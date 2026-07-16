"""
generate_ppt.py
===============
Converts business_presentation.md to business_presentation.pptx
Run: python3 presentations/business/generate_ppt.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ---------------------------------------------------------------------------
# Theme colours
# ---------------------------------------------------------------------------
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT     = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy
HIGHLIGHT  = RGBColor(0x0F, 0x3C, 0x96)   # blue
GOLD       = RGBColor(0xF5, 0xA6, 0x23)   # amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_table(slide, headers, rows, left, top, width, height,
              header_bg=HIGHLIGHT, row_bg=ACCENT, alt_bg=DARK_BG):
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, cols, left, top, width, height).table

    # Column widths equally distributed
    col_w = width // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    def _cell(r, c, text, bg, fg=WHITE, bold=False, sz=12):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg

    for ci, h in enumerate(headers):
        _cell(0, ci, h, header_bg, WHITE, bold=True, sz=13)

    for ri, row in enumerate(rows):
        bg = row_bg if ri % 2 == 0 else alt_bg
        for ci, val in enumerate(row):
            _cell(ri + 1, ci, val, bg, WHITE, sz=11)

    return tbl


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK_BG)

    # Decorative bar
    add_rect(slide, 0, Inches(5.5), SLIDE_W, Inches(0.08), GOLD)

    # Title
    add_textbox(slide, "Protecting Our Network with AI",
                Inches(1), Inches(1.5), Inches(11.33), Inches(1.5),
                font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, "A Machine Learning Intrusion Detection System",
                Inches(1), Inches(3), Inches(11.33), Inches(0.7),
                font_size=22, color=GOLD, align=PP_ALIGN.CENTER)

    # Stats row
    stats = [
        ("99.3%", "Detection Accuracy"),
        ("6×", "Fewer False Alerts"),
        ("<1ms", "Response Time"),
        ("USD 2.15M", "Annual Value"),
    ]
    x_start = Inches(0.5)
    box_w = Inches(2.9)
    for i, (val, lbl) in enumerate(stats):
        x = x_start + i * Inches(3.1)
        add_rect(slide, x, Inches(5.8), box_w, Inches(1.3), HIGHLIGHT)
        add_textbox(slide, val, x, Inches(5.85), box_w, Inches(0.6),
                    font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, x, Inches(6.4), box_w, Inches(0.5),
                    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, "Business Executive Presentation  |  Cybersecurity AI Initiative",
                Inches(1), Inches(7.1), Inches(11.33), Inches(0.4),
                font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


def slide_threat_landscape(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 2", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "Why Traditional Security Is Failing",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    headers = ["Challenge", "Traditional IDS", "Our AI Solution"]
    rows = [
        ["New attack types", "✗  Cannot detect", "✓  Learns patterns"],
        ["Alert volume", "✗  50,000+/day false alerts", "✓  ~7,000/day"],
        ["Response speed", "✗  Minutes", "✓  < 1 second"],
        ["Adaptability", "✗  Manual signature updates", "✓  Retrains on new data"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.4), Inches(12.5), Inches(2.6))

    # Cost callouts
    callouts = [
        ("USD 4.45M", "Average breach cost\n(IBM 2023)"),
        ("207 days", "Mean time to detect\na breach"),
        ("80%", "Of breaches involve\nnetwork access"),
    ]
    for i, (val, lbl) in enumerate(callouts):
        x = Inches(0.4) + i * Inches(4.3)
        add_rect(slide, x, Inches(4.3), Inches(3.9), Inches(1.5), HIGHLIGHT)
        add_textbox(slide, val, x, Inches(4.35), Inches(3.9), Inches(0.65),
                    font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, x, Inches(4.9), Inches(3.9), Inches(0.8),
                    font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


def slide_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 3", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "The Solution in Plain English",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    steps = [
        ("1", "Network\nConnection", "Every packet entering or leaving the network"),
        ("2", "41 Features\nAnalysed", "Speed, size, error rates, connection patterns"),
        ("3", "AI Classifies\nin <1ms", "XGBoost model trained on 125,000+ examples"),
        ("4", "Instant\nDecision", "Normal → Allow  |  Suspicious → Alert SOC"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.3) + i * Inches(3.25)
        add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(2.8), ACCENT)
        add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(0.55), HIGHLIGHT)
        add_textbox(slide, num, x, Inches(1.5), Inches(3.0), Inches(0.55),
                    font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, x, Inches(2.1), Inches(3.0), Inches(0.9),
                    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, x, Inches(3.1), Inches(3.0), Inches(1.0),
                    font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < 3:
            add_textbox(slide, "→", x + Inches(3.05), Inches(2.6), Inches(0.2), Inches(0.5),
                        font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Explainability callout
    add_rect(slide, Inches(0.4), Inches(4.6), Inches(12.5), Inches(1.5), HIGHLIGHT)
    add_textbox(slide, "No Black Boxes — Every Decision is Explainable",
                Inches(0.6), Inches(4.65), Inches(12), Inches(0.45),
                font_size=16, bold=True, color=GOLD)
    add_textbox(slide,
                '"This connection was flagged because it made 511 rapid connection attempts '
                'with zero completed handshakes — a SYN flood signature."',
                Inches(0.6), Inches(5.1), Inches(12), Inches(0.8),
                font_size=13, color=WHITE, italic=True)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 4", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "What the AI Achieves",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    headers = ["Business Outcome", "Our AI", "Traditional IDS"]
    rows = [
        ["Attacks caught", "99.3%", "~70%"],
        ["False alarm rate", "0.7%", "~5%"],
        ["Response time", "< 1ms", "Minutes"],
        ["Zero-day detection", "Yes", "No"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.4), Inches(8.5), Inches(2.5))

    # KPI boxes on right
    kpis = [
        ("993 / 1,000", "Attacks detected\nper 1M connections"),
        ("~7,000", "False alarms/day\nvs. 50,000 traditional"),
        ("4 hrs/day", "SOC analyst time\nsaved per day"),
    ]
    for i, (val, lbl) in enumerate(kpis):
        y = Inches(1.4) + i * Inches(1.1)
        add_rect(slide, Inches(9.2), y, Inches(3.7), Inches(0.95), HIGHLIGHT)
        add_textbox(slide, val, Inches(9.2), y + Inches(0.05), Inches(3.7), Inches(0.45),
                    font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, Inches(9.2), y + Inches(0.45), Inches(3.7), Inches(0.45),
                    font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_roi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 5", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "The Business Case — ROI",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    headers = ["Scenario", "Probability", "Avoided Cost", "Contribution"]
    rows = [
        ["Data breach prevented", "30% annual risk", "USD 4.45M", "USD 1.33M"],
        ["Regulatory fine avoided (GDPR)", "15% risk", "USD 2.0M", "USD 300K"],
        ["Incident response reduction", "80% fewer incidents", "USD 150K", "USD 120K"],
        ["SOC productivity gain", "4 hrs/day × 5 analysts", "USD 400K/yr", "USD 400K"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.4), Inches(12.5), Inches(2.5))

    # Summary boxes
    add_rect(slide, Inches(0.4), Inches(4.2), Inches(5.8), Inches(1.5), HIGHLIGHT)
    add_textbox(slide, "USD 2.15M", Inches(0.4), Inches(4.25), Inches(5.8), Inches(0.7),
                font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Total Estimated Annual Value",
                Inches(0.4), Inches(4.85), Inches(5.8), Inches(0.5),
                font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(6.6), Inches(4.2), Inches(6.3), Inches(1.5), ACCENT)
    add_textbox(slide, "Implementation Cost",
                Inches(6.6), Inches(4.25), Inches(6.3), Inches(0.45),
                font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    cost_lines = "USD 15K one-time deployment   |   USD 20K/yr maintenance\nExisting infrastructure — no new hardware"
    add_textbox(slide, cost_lines,
                Inches(6.6), Inches(4.7), Inches(6.3), Inches(0.8),
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "ROI: >100× in year one",
                Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.6),
                font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_risk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 6", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "What Could Go Wrong — And How We Address It",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)

    headers = ["Risk", "Likelihood", "Impact", "Mitigation"]
    rows = [
        ["Dataset is 25 years old", "Medium", "New attacks not recognized", "Quarterly retraining on modern data"],
        ["Rare attacks missed (U2R)", "Low", "High-severity breach", "Dedicated U2R detection layer"],
        ["Wrong prediction (0.7% FPR)", "Very Low", "False alarm", "Human-in-the-loop review"],
        ["System downtime", "Very Low", "Detection gap", "Fallback to traditional IDS"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.2))

    add_rect(slide, Inches(0.4), Inches(4.9), Inches(12.5), Inches(0.7), HIGHLIGHT)
    add_textbox(slide,
                "Our AI does not replace security analysts — it empowers them.",
                Inches(0.5), Inches(4.95), Inches(12), Inches(0.5),
                font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_ethical_ai(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 7", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "Responsible AI Deployment",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    pillars = [
        ("Explainability", "Every AI decision explained in plain language.\nNo black boxes. Compliant with EU AI Act."),
        ("Fairness Audit", "Tested across TCP, UDP, ICMP traffic.\nNo bias toward any protocol group."),
        ("Privacy", "Operates on metadata only — no payload content.\nFully GDPR compliant."),
        ("Human Oversight", "AI generates alerts;\nanalysts make final decisions."),
        ("Data Security", "Deployed on-premises.\nNo data leaves the network perimeter."),
    ]
    cols = 3
    for i, (title, body) in enumerate(pillars):
        row, col = divmod(i, cols)
        x = Inches(0.3) + col * Inches(4.35)
        y = Inches(1.5) + row * Inches(2.0)
        add_rect(slide, x, y, Inches(4.1), Inches(1.8), ACCENT)
        add_rect(slide, x, y, Inches(4.1), Inches(0.5), HIGHLIGHT)
        add_textbox(slide, title, x, y + Inches(0.05), Inches(4.1), Inches(0.45),
                    font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, body, x, y + Inches(0.55), Inches(4.1), Inches(1.1),
                    font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


def slide_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 8", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "Why AI Wins",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    headers = ["Traditional Signature IDS", "Our AI-Based IDS"]
    rows = [
        ["✗  Misses all new / unknown attacks", "✓  Detects unknown attacks by behaviour"],
        ["✗  Requires daily manual signature updates", "✓  Self-improving — retrains on new data"],
        ["✗  50,000+ daily alerts overwhelm SOC", "✓  ~7,000 targeted high-confidence alerts"],
        ["✗  Static — no learning over time", "✓  Continuously improves with retraining"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.4), Inches(12.5), Inches(2.8))

    # Case study
    add_rect(slide, Inches(0.4), Inches(4.4), Inches(12.5), Inches(2.4), ACCENT)
    add_textbox(slide, "Case Study — SYN Flood Detection",
                Inches(0.6), Inches(4.45), Inches(12), Inches(0.5),
                font_size=16, bold=True, color=GOLD)
    add_textbox(slide,
                "Signature IDS:  Catches only known SYN flood patterns — misses novel variants.\n"
                "Our AI:  Detected 100% of SYN floods in testing, including novel variants,\n"
                "          by learning serror_rate + count feature combinations.",
                Inches(0.6), Inches(4.95), Inches(12), Inches(1.6),
                font_size=13, color=WHITE)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 9", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "Phased Rollout Plan",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    phases = [
        ("Phase 1", "Weeks 1–4", "Pilot Deployment",
         "• Deploy on 1 network segment\n• Run in shadow mode (alert, don't block)\n• Validate against current IDS"),
        ("Phase 2", "Weeks 5–12", "Validation",
         "• Compare AI alerts vs. analyst verdicts\n• Fine-tune threshold for risk tolerance\n• Train SOC team on AI dashboard"),
        ("Phase 3", "Month 4+", "Full Deployment",
         "• Replace signature IDS as primary system\n• Quarterly retraining with new data\n• Integrate with SIEM (Splunk / Sentinel)"),
    ]
    for i, (phase, timing, title, bullets) in enumerate(phases):
        x = Inches(0.3) + i * Inches(4.35)
        add_rect(slide, x, Inches(1.5), Inches(4.1), Inches(4.5), ACCENT)
        add_rect(slide, x, Inches(1.5), Inches(4.1), Inches(0.55), HIGHLIGHT)
        add_textbox(slide, f"{phase}  |  {timing}", x, Inches(1.52), Inches(4.1), Inches(0.5),
                    font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, x, Inches(2.1), Inches(4.1), Inches(0.55),
                    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, bullets, x, Inches(2.75), Inches(4.1), Inches(2.8),
                    font_size=13, color=LIGHT_GREY)

    add_textbox(slide, "Minimal disruption — existing infrastructure, no hardware changes required.",
                Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.45),
                font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), GOLD)

    add_textbox(slide, "SLIDE 10", Inches(0.4), Inches(0.2), Inches(3), Inches(0.4),
                font_size=11, color=GOLD, bold=True)
    add_textbox(slide, "Our Recommendation: Proceed",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
                font_size=30, bold=True, color=WHITE)

    # Decision box
    add_rect(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(1.1), HIGHLIGHT)
    add_textbox(slide, "Approve Phase 1 Pilot Budget:  USD 15,000",
                Inches(0.5), Inches(1.45), Inches(12), Inches(0.5),
                font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "One-time deployment on the DMZ network segment",
                Inches(0.5), Inches(1.85), Inches(12), Inches(0.45),
                font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # Why now
    add_textbox(slide, "Why Now:",
                Inches(0.4), Inches(2.75), Inches(3), Inches(0.45),
                font_size=16, bold=True, color=GOLD)
    why = ("• Cyber threats growing 38% year-over-year\n"
           "• AI technology proven — 99.3% F1 on industry benchmark\n"
           "• Implementation cost is low relative to risk exposure")
    add_textbox(slide, why, Inches(0.4), Inches(3.2), Inches(6.0), Inches(1.3),
                font_size=13, color=WHITE)

    # Next steps
    add_textbox(slide, "Immediate Next Steps:",
                Inches(6.8), Inches(2.75), Inches(6), Inches(0.45),
                font_size=16, bold=True, color=GOLD)
    steps = ("✓  Approve pilot on DMZ network segment\n"
             "✓  Assign 1 SOC analyst as AI system owner\n"
             "✓  Schedule retraining with CICIDS-2017 dataset\n"
             "✓  Define escalation policy for high-priority alerts")
    add_textbox(slide, steps, Inches(6.8), Inches(3.2), Inches(6.2), Inches(1.3),
                font_size=13, color=WHITE)

    add_rect(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.7), ACCENT)
    add_textbox(slide,
                "Live demo available on request  |  Full technical report & bias audit available from the project team",
                Inches(0.5), Inches(5.65), Inches(12), Inches(0.5),
                font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    add_textbox(slide, "Questions?",
                Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.6),
                font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_threat_landscape(prs)
    slide_how_it_works(prs)
    slide_results(prs)
    slide_roi(prs)
    slide_risk(prs)
    slide_ethical_ai(prs)
    slide_comparison(prs)
    slide_roadmap(prs)
    slide_next_steps(prs)

    out = Path(__file__).parent / "business_presentation.pptx"
    prs.save(str(out))
    print(f"✅  Saved: {out}")


if __name__ == "__main__":
    main()
