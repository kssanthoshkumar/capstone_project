"""
generate_ppt.py
===============
Converts technical_presentation.md to technical_presentation.pptx
Run: python3 presentations/technical/generate_ppt.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep blue-black
ACCENT     = RGBColor(0x1B, 0x2A, 0x3D)   # panel bg
HIGHLIGHT  = RGBColor(0x0E, 0x4D, 0x92)   # corporate blue
TEAL       = RGBColor(0x00, 0x8B, 0x8B)   # teal accent
GOLD       = RGBColor(0xF5, 0xA6, 0x23)   # amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
CODE_BG    = RGBColor(0x10, 0x10, 0x1A)   # code block bg

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_table(slide, headers, rows, left, top, width, height,
              header_bg=HIGHLIGHT, row_bg=ACCENT, alt_bg=DARK_BG, font_size=12):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows) + 1, cols, left, top, width, height).table
    col_w = width // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    def _cell(r, c, text, bg, fg=WHITE, bold=False, sz=12):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg

    for ci, h in enumerate(headers):
        _cell(0, ci, h, header_bg, WHITE, bold=True, sz=font_size)
    for ri, row in enumerate(rows):
        bg = row_bg if ri % 2 == 0 else alt_bg
        for ci, val in enumerate(row):
            _cell(ri + 1, ci, val, bg, WHITE, sz=font_size - 1)

    return tbl


def slide_label(slide, num):
    add_textbox(slide, f"SLIDE {num}", Inches(0.4), Inches(0.18),
                Inches(3), Inches(0.38), font_size=11, color=TEAL, bold=True)


def divider(slide):
    add_rect(slide, 0, Inches(1.12), SLIDE_W, Inches(0.05), GOLD)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(5.4), SLIDE_W, Inches(0.06), GOLD)

    add_textbox(slide, "Network Traffic Anomaly Detection",
                Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.2),
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "using Machine Learning",
                Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.7),
                font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Detecting Cyber Threats with AI: From Data to Deployment",
                Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.55),
                font_size=18, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    tags = [
        ("Domain", "Cybersecurity"),
        ("Dataset", "NSL-KDD"),
        ("Models", "LR · RF · XGBoost · IsolForest · Autoencoder"),
        ("Best Result", "Autoencoder F1=0.836 (recall)  |  XGBoost AUC=0.967 (precision)"),
    ]
    for i, (k, v) in enumerate(tags):
        x = Inches(0.3) + i * Inches(3.25)
        add_rect(slide, x, Inches(5.6), Inches(3.1), Inches(1.2), ACCENT)
        add_textbox(slide, k, x, Inches(5.65), Inches(3.1), Inches(0.4),
                    font_size=11, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, v, x, Inches(6.0), Inches(3.1), Inches(0.65),
                    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Capstone Project  |  Pillar 5  |  AI/ML Fundamentals",
                Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.35),
                font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


def slide_problem_framing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 2)
    add_textbox(slide, "Problem Framing & Task Type",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    # Task mapping panel
    add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.4), CODE_BG)
    add_textbox(slide, "Task Mapping",
                Inches(0.5), Inches(1.42), Inches(5.8), Inches(0.42),
                font_size=13, bold=True, color=TEAL)
    task_text = (
        "Binary Classification (primary)\n"
        "  →  Normal vs. Attack\n\n"
        "Multi-class (secondary)\n"
        "  →  DoS | Probe | R2L | U2R\n\n"
        "Unsupervised (benchmark)\n"
        "  →  Isolation Forest + Autoencoder"
    )
    add_textbox(slide, task_text, Inches(0.5), Inches(1.85), Inches(5.7), Inches(1.8),
                font_size=13, color=LIGHT_GREY)

    # Problem statement
    add_rect(slide, Inches(6.6), Inches(1.4), Inches(6.4), Inches(1.0), HIGHLIGHT)
    add_textbox(slide, "Problem Statement",
                Inches(6.7), Inches(1.43), Inches(6.2), Inches(0.38),
                font_size=12, bold=True, color=GOLD)
    add_textbox(slide,
                "Classify network connection records\n(41 features) as Normal or Attack,\nidentifying attack category.",
                Inches(6.7), Inches(1.8), Inches(6.2), Inches(0.55),
                font_size=13, color=WHITE)

    # Metrics table
    add_textbox(slide, "Success Metrics",
                Inches(6.6), Inches(2.55), Inches(6), Inches(0.4),
                font_size=14, bold=True, color=TEAL)
    headers = ["Metric", "Target", "Rationale"]
    rows = [
        ["F1-Score", "≥ 0.97", "Imbalanced classes"],
        ["AUC-ROC", "≥ 0.99", "Threshold-agnostic"],
        ["FPR", "≤ 0.02", "SOC capacity"],
        ["Recall", "≥ 0.98", "Attack miss cost"],
    ]
    add_table(slide, headers, rows,
              Inches(6.6), Inches(2.95), Inches(6.4), Inches(2.1), font_size=13)


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 3)
    add_textbox(slide, "Dataset Overview — NSL-KDD",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    headers = ["Property", "Value"]
    rows = [
        ["Train rows", "125,973"],
        ["Test rows", "22,544"],
        ["Features", "41  (38 numeric + 3 categorical)"],
        ["Missing values", "None"],
        ["Attack split", "DoS 76%  |  Probe 19%  |  R2L 4%  |  U2R 0.3%"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.4), Inches(6.2), Inches(2.4), font_size=13)

    add_rect(slide, Inches(0.4), Inches(4.0), Inches(6.2), Inches(0.65), ACCENT)
    add_textbox(slide,
                "Key advantage: removed 78% duplicate records vs KDD'99 — avoids inflated metrics.",
                Inches(0.5), Inches(4.05), Inches(6.0), Inches(0.55),
                font_size=12, color=GOLD, italic=True)

    # Attack categories
    attacks = [
        ("DoS", "neptune, smurf", "Overwhelm service availability"),
        ("Probe", "ipsweep, nmap", "Network reconnaissance/scanning"),
        ("R2L", "guess_passwd", "Remote-to-local unauthorised access"),
        ("U2R", "buffer_overflow", "Privilege escalation to root"),
    ]
    add_textbox(slide, "Attack Categories",
                Inches(7.0), Inches(1.35), Inches(6), Inches(0.42),
                font_size=14, bold=True, color=TEAL)
    for i, (cat, examples, desc) in enumerate(attacks):
        y = Inches(1.85) + i * Inches(1.1)
        add_rect(slide, Inches(7.0), y, Inches(6.0), Inches(1.0), ACCENT)
        add_rect(slide, Inches(7.0), y, Inches(1.2), Inches(1.0), HIGHLIGHT)
        add_textbox(slide, cat, Inches(7.0), y + Inches(0.25),
                    Inches(1.2), Inches(0.5),
                    font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, examples, Inches(8.3), y + Inches(0.05),
                    Inches(4.6), Inches(0.38), font_size=12, color=TEAL, bold=True)
        add_textbox(slide, desc, Inches(8.3), y + Inches(0.45),
                    Inches(4.6), Inches(0.45), font_size=12, color=LIGHT_GREY)


def slide_preprocessing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 4)
    add_textbox(slide, "Preprocessing Pipeline",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    steps = [
        ("Remove Duplicates", "No duplicate rows remain"),
        ("Outlier Capping", "IQR × 3 — computed on train only"),
        ("Numeric (38)", "Median Imputation → StandardScaler"),
        ("Categorical (3)", "Mode Imputation → OneHotEncoder"),
        ("Feature Eng.", "+10 domain features (byte_ratio, error_rate_total …)"),
        ("Binning", "count, srv_count → low/med/high/very_high"),
        ("Output", "55-feature scaled matrix  |  125k × 55 train"),
    ]
    for i, (step, detail) in enumerate(steps):
        y = Inches(1.35) + i * Inches(0.73)
        add_rect(slide, Inches(0.4), y, Inches(2.6), Inches(0.63), HIGHLIGHT)
        add_textbox(slide, step, Inches(0.45), y + Inches(0.1),
                    Inches(2.5), Inches(0.43), font_size=13, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(3.1), y, Inches(9.8), Inches(0.63), ACCENT)
        add_textbox(slide, detail, Inches(3.2), y + Inches(0.1),
                    Inches(9.6), Inches(0.43), font_size=13, color=LIGHT_GREY)

        if i < len(steps) - 1:
            add_textbox(slide, "↓", Inches(1.5), y + Inches(0.63),
                        Inches(0.5), Inches(0.1), font_size=10,
                        color=GOLD, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5), CODE_BG)
    add_textbox(slide,
                "Anti-leakage:  difficulty_level excluded  |  preprocessor fit on train only  |  outlier caps from train quantiles",
                Inches(0.5), Inches(6.73), Inches(12.2), Inches(0.42),
                font_size=12, color=TEAL, align=PP_ALIGN.CENTER)


def slide_eda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 5)
    add_textbox(slide, "EDA — Key Findings",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    findings = [
        ("Finding 1", "serror_rate Separates Attacks Perfectly",
         "The SYN error rate is near 0 for all normal traffic and near 1 for DoS/SYN-flood attacks.\nSingle strongest binary discriminator in the dataset."),
        ("Finding 2", "flag=S0 ≈ 100% Attack",
         "Connections with flag S0 (incomplete handshake) are overwhelmingly attacks.\nOne-hot encoded as a top-ranked feature."),
        ("Finding 3", "Multicollinearity in Error Features",
         "serror_rate ↔ dst_host_serror_rate: r=0.92\nPCA applied before Autoencoder to decorrelate."),
        ("Finding 4", "t-SNE: Clear Cluster Separation",
         "Attack and normal traffic form well-separated clusters in 2D embedding,\nconfirming that linear + tree models should achieve high accuracy."),
    ]
    for i, (tag, title, body) in enumerate(findings):
        row, col = divmod(i, 2)
        x = Inches(0.4) + col * Inches(6.5)
        y = Inches(1.35) + row * Inches(2.4)
        add_rect(slide, x, y, Inches(6.2), Inches(2.2), ACCENT)
        add_rect(slide, x, y, Inches(6.2), Inches(0.45), HIGHLIGHT)
        add_textbox(slide, tag, x, y + Inches(0.04), Inches(6.2), Inches(0.38),
                    font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, x, y + Inches(0.5), Inches(6.2), Inches(0.45),
                    font_size=14, bold=True, color=WHITE)
        add_textbox(slide, body, x, y + Inches(0.95), Inches(6.2), Inches(1.15),
                    font_size=12, color=LIGHT_GREY)


def slide_feature_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 6)
    add_textbox(slide, "Feature Engineering & Selection",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    # Domain features table
    add_textbox(slide, "10 Domain-Derived Features",
                Inches(0.4), Inches(1.35), Inches(6.2), Inches(0.42),
                font_size=14, bold=True, color=TEAL)
    headers = ["Feature", "Security Insight"]
    rows = [
        ["byte_ratio", "Asymmetric traffic → C2 communication"],
        ["error_rate_total", "High = DoS or scan"],
        ["log_src_bytes", "Normalises heavy-tail distribution"],
        ["is_long_connection", "Tunnelling detection"],
        ["host_srv_ratio", "Port scan concentration"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.8), Inches(6.2), Inches(2.3), font_size=12)

    # Selection methods
    methods = [
        ("Filter\n(Mutual Information)", "20 features selected\nserror_rate, flag_SF,\nflag_S0, logged_in …"),
        ("Embedded\n(Gradient Boosting)", "18 features\nat median threshold"),
        ("Agreement\n(Both Methods)", "15 high-confidence\nfeatures used"),
    ]
    add_textbox(slide, "Feature Selection",
                Inches(6.8), Inches(1.35), Inches(6), Inches(0.42),
                font_size=14, bold=True, color=TEAL)
    for i, (method, result) in enumerate(methods):
        x = Inches(6.8) + i * Inches(2.15)
        add_rect(slide, x, Inches(1.8), Inches(2.0), Inches(2.3), ACCENT)
        add_rect(slide, x, Inches(1.8), Inches(2.0), Inches(0.55), HIGHLIGHT)
        add_textbox(slide, method, x, Inches(1.82), Inches(2.0), Inches(0.5),
                    font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, result, x, Inches(2.4), Inches(2.0), Inches(1.6),
                    font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # PCA
    add_rect(slide, Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.85), CODE_BG)
    add_textbox(slide, "Dimensionality Reduction  —  PCA",
                Inches(0.5), Inches(4.33), Inches(5), Inches(0.38),
                font_size=13, bold=True, color=TEAL)
    add_textbox(slide,
                "55 → 23 components  (95% variance retained)   "
                "|   Used for Autoencoder training   "
                "|   Supervised models use full 55-feature set",
                Inches(0.5), Inches(4.68), Inches(12.2), Inches(0.38),
                font_size=12, color=LIGHT_GREY)


def slide_model_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 7)
    add_textbox(slide, "Model Architectures",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    models = [
        ("Logistic Regression\n(baseline)",
         "L2 regularisation\nC=1.0, lbfgs solver"),
        ("Decision Tree",
         "max_depth=10\nGini impurity\nFast, interpretable"),
        ("Random Forest",
         "200 trees\nStratifiedKFold CV\nmax_depth tuned"),
        ("XGBoost\n(tuned)",
         "300 estimators, depth=6\nscale_pos_weight\nt=0.05 threshold"),
        ("LightGBM",
         "300 estimators\nHistogram-based\nFaster than XGBoost"),
        ("SVM\n(LinearSVC)",
         "CalibratedClassifierCV\nLinear kernel O(n)\nclass_weight='balanced'"),
        ("Isolation Forest",
         "200 trees\ncontamination=0.47\nUnsupervised"),
        ("Autoencoder\n(Deep Learning)",
         "55→64→32→16→32→64→55\nTrained on NORMAL only\nMSE > 95th pct = attack"),
    ]
    n_models = len(models)
    col_w = Inches(12.5 / n_models)
    for i, (name, params) in enumerate(models):
        x = Inches(0.3) + i * (col_w + Inches(0.05))
        highlight = GOLD if "tuned" in name.lower() or "Deep" in name else HIGHLIGHT
        add_rect(slide, x, Inches(1.35), col_w, Inches(5.5), ACCENT)
        add_rect(slide, x, Inches(1.35), col_w, Inches(0.55), highlight)
        add_textbox(slide, name, x, Inches(1.37), col_w, Inches(0.52),
                    font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, params, x, Inches(1.95), col_w, Inches(4.5),
                    font_size=10, color=LIGHT_GREY)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 8)
    add_textbox(slide, "Results Comparison",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    add_textbox(slide, "Results on KDDTest+ (22,544 held-out records with novel attack subtypes)",
                Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.38),
                font_size=12, color=LIGHT_GREY, italic=True)

    headers = ["Model", "Precision", "Recall", "F1", "AUC-ROC", "Type"]
    rows = [
        ["Autoencoder  ★", "0.741", "0.960", "0.836", "0.817", "Deep Learning"],
        ["XGBoost (t=0.05)  ★", "0.966", "0.725", "0.828", "0.967", "Supervised"],
        ["Decision Tree", "0.968", "0.669", "0.791", "0.838", "Supervised"],
        ["XGBoost (default)", "0.967", "0.648", "0.776", "0.967", "Supervised"],
        ["LightGBM", "0.966", "0.631", "0.763", "0.955", "Supervised"],
        ["Random Forest", "0.968", "0.605", "0.744", "0.953", "Supervised"],
        ["Isolation Forest", "0.748", "0.716", "0.732", "0.779", "Unsupervised"],
        ["SVM (LinearSVC)", "0.735", "0.619", "0.672", "0.651", "Supervised"],
        ["Logistic Regression", "0.734", "0.620", "0.673", "0.654", "Supervised"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.65), Inches(12.5), Inches(3.3), font_size=11)

    tradeoffs = [
        ("Highest Recall", "Autoencoder", "0.960"),
        ("Best F1", "Autoencoder", "0.836"),
        ("Best Precision", "XGBoost", "0.967"),
        ("Best AUC", "XGBoost", "0.967"),
    ]
    add_textbox(slide, "Deployment trade-offs:",
                Inches(0.4), Inches(5.15), Inches(4), Inches(0.4),
                font_size=13, bold=True, color=TEAL)
    for i, (priority, model, score) in enumerate(tradeoffs):
        x = Inches(0.4) + i * Inches(3.1)
        add_rect(slide, x, Inches(5.6), Inches(2.9), Inches(1.1), ACCENT)
        add_textbox(slide, priority, x, Inches(5.62), Inches(2.9), Inches(0.38),
                    font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"{model}: {score}", x, Inches(5.98), Inches(2.9), Inches(0.35),
                    font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_shap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 9)
    add_textbox(slide, "SHAP + LIME + PDP + ICE — Explainability Suite",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    features = [
        ("1", "serror_rate",
         "High → DoS / SYN flood\nPushes prediction strongly to 'attack'"),
        ("2", "same_srv_rate",
         "Low → port scan\nDistinguishes probe from normal"),
        ("3", "flag_S0",
         "Incomplete connection = attack signature\nOne-hot encoded top feature"),
        ("4", "dst_host_serror_rate",
         "Persistent SYN errors at destination\nCorrelated with serror_rate"),
        ("5", "src_bytes",
         "Zero or extreme = anomalous\nKey DoS indicator"),
    ]
    for i, (rank, feat, insight) in enumerate(features):
        y = Inches(1.35) + i * Inches(0.96)
        add_rect(slide, Inches(0.4), y, Inches(0.6), Inches(0.86), HIGHLIGHT)
        add_textbox(slide, rank, Inches(0.4), y + Inches(0.2),
                    Inches(0.6), Inches(0.45), font_size=18, bold=True,
                    color=GOLD, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.1), y, Inches(12.0), Inches(0.86), ACCENT)
        add_textbox(slide, feat, Inches(1.2), y + Inches(0.03),
                    Inches(4.5), Inches(0.42), font_size=15, bold=True, color=WHITE)
        add_textbox(slide, insight, Inches(1.2), y + Inches(0.44),
                    Inches(11.6), Inches(0.38), font_size=12, color=LIGHT_GREY)

    add_rect(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.85), CODE_BG)
    add_textbox(slide, "Local explanation — neptune DoS instance #7841:",
                Inches(0.5), Inches(6.33), Inches(6), Inches(0.38),
                font_size=12, bold=True, color=TEAL)
    add_textbox(slide,
                "serror_rate=1.0 → +1.42 SHAP   |   flag=S0 → +0.87 SHAP   |   count=511 → +0.63 SHAP",
                Inches(0.5), Inches(6.68), Inches(12.2), Inches(0.38),
                font_size=12, color=GOLD)


def slide_bias_audit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 10)
    add_textbox(slide, "Bias Audit Results",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    add_textbox(slide, "Fairness evaluated across 3 subgroup dimensions",
                Inches(0.4), Inches(1.2), Inches(12), Inches(0.42),
                font_size=15, color=LIGHT_GREY, italic=True)

    headers = ["Subgroup", "DPD", "EOD_FPR", "EOD_TPR", "FPR_Ratio", "Status"]
    rows = [
        ["protocol_type", "0.031", "0.007", "0.013", "1.19", "✅ PASS"],
        ["service", "0.044", "0.018", "0.021", "1.22", "✅ PASS"],
        ["flag", "0.052", "0.021", "0.028", "1.24", "✅ PASS"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.7), Inches(12.5), Inches(1.8), font_size=13)

    thresholds = [
        ("DPD > 0.05", "Concern threshold"),
        ("EOD > 0.10", "Concern threshold"),
        ("Ratio > 1.25", "Concern threshold"),
    ]
    add_textbox(slide, "Concern Thresholds (all subgroups below):",
                Inches(0.4), Inches(3.65), Inches(6), Inches(0.42),
                font_size=14, bold=True, color=TEAL)
    for i, (t, lbl) in enumerate(thresholds):
        x = Inches(0.4) + i * Inches(4.2)
        add_rect(slide, x, Inches(4.1), Inches(3.9), Inches(0.9), ACCENT)
        add_textbox(slide, t, x, Inches(4.12), Inches(3.9), Inches(0.42),
                    font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, x, Inches(4.52), Inches(3.9), Inches(0.38),
                    font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.0), CODE_BG)
    add_textbox(slide,
                "Slight variation in ICMP/UDP reflects genuine traffic pattern differences,\n"
                "not model bias. All subgroups pass fairness thresholds.",
                Inches(0.5), Inches(5.25), Inches(12.2), Inches(0.85),
                font_size=13, color=LIGHT_GREY, italic=True)


def slide_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 11)
    add_textbox(slide, "Limitations & Mitigations",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    headers = ["Limitation", "Impact", "Mitigation"]
    rows = [
        ["Class Imbalance (U2R: 0.3%)", "Misses rare high-severity attacks", "SMOTE + class_weight='balanced'"],
        ["Dataset Age (1999 traffic)", "Zero-day attacks not represented", "Retrain on CICIDS-2017/2018"],
        ["Overfitting (train F1=0.999)", "Marginal; CV-controlled", "StratifiedKFold + depth limits"],
        ["Data Leakage risk", "Eliminated in pipeline", "difficulty_level excluded; train-only fit"],
        ["Distribution Shift", "Production degradation over time", "Evidently AI drift monitoring"],
    ]
    add_table(slide, headers, rows,
              Inches(0.4), Inches(1.35), Inches(12.5), Inches(3.5), font_size=12)

    add_rect(slide, Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.85), HIGHLIGHT)
    add_textbox(slide,
                "Recommended next step: retrain quarterly on CICIDS-2017/2018 modern captures to reduce dataset-age risk.",
                Inches(0.5), Inches(5.1), Inches(12.2), Inches(0.7),
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    divider(slide)
    slide_label(slide, 12)
    add_textbox(slide, "Deployment Architecture & Conclusion",
                Inches(0.4), Inches(0.5), Inches(12), Inches(0.65),
                font_size=28, bold=True, color=WHITE)

    # Architecture diagram — use the PNG if it exists, fallback to text
    arch_img = Path(__file__).parent.parent / "architecture_diagram.png"
    if arch_img.exists():
        # Image aspect ratio ≈ 20:12; fit into left 7.8" × 4.7"
        slide.shapes.add_picture(
            str(arch_img),
            Inches(0.3), Inches(1.3),
            Inches(7.8), Inches(4.68),
        )
    else:
        # Fallback: text representation
        arch_lines = [
            ("Client Request", LIGHT_GREY),
            ("POST /predict  →  FastAPI  (app.py)", WHITE),
            ("Pydantic validation  (41 features, range checks)", LIGHT_GREY),
            ("preprocessor.pkl  →  ColumnTransformer", LIGHT_GREY),
            ("xgboost.pkl  →  predict_proba()", LIGHT_GREY),
            ("Response:  { prediction, label, probability }", GOLD),
        ]
        add_rect(slide, Inches(0.4), Inches(1.35), Inches(7.0), Inches(4.5), CODE_BG)
        for i, (line, color) in enumerate(arch_lines):
            add_textbox(slide, line, Inches(0.5), Inches(1.45) + i * Inches(0.6),
                        Inches(6.8), Inches(0.5), font_size=13, color=color)

    # Summary bullets (right side)
    summary = [
        ("F1=0.776 XGB  |  F1=0.836 AE", "All models evaluated on KDDTest+"),
        ("SHAP + LIME + PDP + ICE", "Domain-meaningful features confirmed"),
        ("Bias Audit (DIR=0.21)", "Protocol-level disparate impact explained"),
        ("FastAPI + Streamlit", "Input validation + X-From traceability"),
        ("Full Reproducibility", "Saved models, configs & notebooks"),
    ]
    add_textbox(slide, "Summary",
                Inches(8.3), Inches(1.35), Inches(4.8), Inches(0.42),
                font_size=16, bold=True, color=TEAL)
    for i, (title, detail) in enumerate(summary):
        y = Inches(1.85) + i * Inches(0.88)
        add_rect(slide, Inches(8.3), y, Inches(4.8), Inches(0.78), ACCENT)
        add_textbox(slide, title, Inches(8.4), y + Inches(0.03),
                    Inches(4.6), Inches(0.38), font_size=12, bold=True, color=GOLD)
        add_textbox(slide, detail, Inches(8.4), y + Inches(0.4),
                    Inches(4.6), Inches(0.33), font_size=11, color=LIGHT_GREY)

    add_rect(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.48), HIGHLIGHT)
    add_textbox(slide,
                "Technical Presentation  |  Capstone Project  |  AI/ML Fundamentals",
                Inches(0.5), Inches(6.78), Inches(12.2), Inches(0.38),
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER, italic=True)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem_framing(prs)
    slide_dataset(prs)
    slide_preprocessing(prs)
    slide_eda(prs)
    slide_feature_engineering(prs)
    slide_model_arch(prs)
    slide_results(prs)
    slide_shap(prs)
    slide_bias_audit(prs)
    slide_limitations(prs)
    slide_deployment(prs)

    out = Path(__file__).parent / "technical_presentation.pptx"
    prs.save(str(out))
    print(f"✅  Saved: {out}")


if __name__ == "__main__":
    main()
